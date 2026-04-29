"""
Generate Tax Lens presentation deck.
Run: pip install python-pptx && python generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Brand palette
INDIGO = RGBColor(0x1D, 0x4E, 0xD8)
PURPLE = RGBColor(0x63, 0x66, 0xF1)
SLATE  = RGBColor(0x1F, 0x29, 0x37)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
LIGHT  = RGBColor(0xF4, 0xF6, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x04, 0x78, 0x57)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_band(slide, y, h, color):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, prs.slide_width, h)
    band.fill.solid(); band.fill.fore_color.rgb = color
    band.line.fill.background()
    return band


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=14, color=SLATE, bullet_color=INDIGO):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet marker
        bm = p.add_run()
        bm.text = "■  "
        bm.font.color.rgb = bullet_color
        bm.font.size = Pt(size)
        bm.font.bold = True
        # text
        r = p.add_run()
        r.text = item
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        p.space_after = Pt(6)


def add_card(slide, x, y, w, h, title, body_lines, accent=INDIGO):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB); card.line.width = Pt(0.75)
    # accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    stripe.fill.solid(); stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    add_text(slide, x + Inches(0.25), y + Inches(0.12), w - Inches(0.4), Inches(0.4),
             title, size=14, bold=True, color=accent)
    add_bullets(slide, x + Inches(0.25), y + Inches(0.55), w - Inches(0.4), h - Inches(0.6),
                body_lines, size=11, color=SLATE, bullet_color=accent)


def add_footer(slide, page_num, total):
    add_band(slide, prs.slide_height - Inches(0.4), Inches(0.4), LIGHT)
    add_text(slide, Inches(0.5), prs.slide_height - Inches(0.35), Inches(8), Inches(0.3),
             "Tax Lens · Codeathon 2026", size=10, color=GRAY)
    add_text(slide, prs.slide_width - Inches(2.5), prs.slide_height - Inches(0.35),
             Inches(2), Inches(0.3),
             f"{page_num} / {total}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# ==============================================================================
# SLIDE 1 — TITLE
# ==============================================================================
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, color=WHITE)
    # Big diagonal banner
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(7.5))
    band.fill.solid(); band.fill.fore_color.rgb = INDIGO
    band.line.fill.background()
    band2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(7), 0, Inches(7), Inches(7.5))
    band2.fill.solid(); band2.fill.fore_color.rgb = PURPLE
    band2.line.fill.background()

    add_text(s, Inches(0.8), Inches(2.2), Inches(11), Inches(1),
             "TAX LENS", size=72, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.4), Inches(11), Inches(0.6),
             "Natural-Language Tax Advisory · 2024 · Azure OpenAI Agentic Chat",
             size=22, color=WHITE)
    add_text(s, Inches(0.8), Inches(4.6), Inches(11), Inches(0.5),
             "Angular 17 + PrimeNG  ·  FastAPI + Python  ·  MongoDB  ·  GPT Function Calling",
             size=16, color=WHITE)
    add_text(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
             "Codeathon Presentation · Architecture & Run Guide",
             size=12, color=WHITE)

# ==============================================================================
# SLIDE 2 — PROBLEM / SOLUTION
# ==============================================================================
def slide_problem():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, LIGHT)
    add_band(s, 0, Inches(0.9), INDIGO)
    add_text(s, Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
             "The Problem & Our Solution", size=28, bold=True, color=WHITE)

    add_card(s, Inches(0.5), Inches(1.2), Inches(6), Inches(5.6),
             "PROBLEM",
             [
                 "Tax software shows numbers — but not WHY.",
                 "Users don't trust black-box calculations.",
                 "They want to ask: 'What if we moved to Texas?'",
                 "They want to challenge: 'Why is our rate so high?'",
                 "They want to compare companies in plain English.",
                 "AI-only tools hallucinate tax math — dangerous.",
             ], accent=ORANGE)

    add_card(s, Inches(6.8), Inches(1.2), Inches(6), Inches(5.6),
             "SOLUTION",
             [
                 "Deterministic 7-step rule-based tax engine (Python).",
                 "AI agent layer with 13 callable tools (Azure OpenAI).",
                 "LLM never computes — only orchestrates and explains.",
                 "Persistent session memory in MongoDB → real continuity.",
                 "Interactive grid + chat UI in Angular + PrimeNG.",
                 "Every number is auditable, traceable, defensible.",
             ], accent=GREEN)

    add_footer(s, 2, TOTAL)

# ==============================================================================
# SLIDE 3 — ARCHITECTURE
# ==============================================================================
def slide_architecture():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, WHITE)
    add_band(s, 0, Inches(0.9), INDIGO)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "Architecture · 3 layers, 1 source of truth", size=28, bold=True, color=WHITE)

    # Frontend block
    fe = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(3.6), Inches(4.8))
    fe.fill.solid(); fe.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xFF)
    fe.line.color.rgb = INDIGO
    add_text(s, Inches(0.7), Inches(1.55), Inches(3.4), Inches(0.4),
             "FRONTEND", size=14, bold=True, color=INDIGO)
    add_text(s, Inches(0.7), Inches(1.95), Inches(3.4), Inches(0.5),
             "Angular 17 · PrimeNG", size=18, bold=True, color=SLATE)
    add_bullets(s, Inches(0.8), Inches(2.6), Inches(3.4), Inches(3.5), [
        "Standalone components",
        "PrimeNG p-table grid",
        "Reactive signals state",
        "Conversational chat panel",
        "Tool-call chips (agent transparency)",
        "Step-by-step breakdown card",
    ], size=12)

    # Backend block
    be = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.85), Inches(1.4), Inches(3.6), Inches(4.8))
    be.fill.solid(); be.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
    be.line.color.rgb = GREEN
    add_text(s, Inches(4.95), Inches(1.55), Inches(3.4), Inches(0.4),
             "BACKEND", size=14, bold=True, color=GREEN)
    add_text(s, Inches(4.95), Inches(1.95), Inches(3.4), Inches(0.5),
             "FastAPI · Python", size=18, bold=True, color=SLATE)
    add_bullets(s, Inches(5.05), Inches(2.6), Inches(3.4), Inches(3.5), [
        "Excel parser (pandas)",
        "Pure 7-step tax engine",
        "13 tool functions",
        "REST endpoints",
        "MongoDB session store",
        "In-memory pre-computed cache",
    ], size=12, bullet_color=GREEN)

    # AI block
    ai = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.1), Inches(1.4), Inches(3.6), Inches(4.8))
    ai.fill.solid(); ai.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    ai.line.color.rgb = ORANGE
    add_text(s, Inches(9.2), Inches(1.55), Inches(3.4), Inches(0.4),
             "AI LAYER", size=14, bold=True, color=ORANGE)
    add_text(s, Inches(9.2), Inches(1.95), Inches(3.4), Inches(0.5),
             "Azure OpenAI", size=18, bold=True, color=SLATE)
    add_bullets(s, Inches(9.3), Inches(2.6), Inches(3.4), Inches(3.5), [
        "Function calling (agentic)",
        "13 tool schemas",
        "Multi-iteration tool loop",
        "Persona + tone prompt",
        "Memory via message history",
        "NEVER computes — only narrates",
    ], size=12, bullet_color=ORANGE)

    # Arrows
    for x_start, x_end in [(Inches(4.2), Inches(4.85)), (Inches(8.45), Inches(9.1))]:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x_start, Inches(3.65), x_end - x_start, Inches(0.4))
        ar.fill.solid(); ar.fill.fore_color.rgb = GRAY
        ar.line.fill.background()

    add_text(s, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
             "Key principle: numbers are deterministic; only narration is AI.",
             size=14, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

    add_footer(s, 3, TOTAL)

# ==============================================================================
# SLIDE 4 — TAX ENGINE 7 STEPS
# ==============================================================================
def slide_tax_steps():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, WHITE)
    add_band(s, 0, Inches(0.9), INDIGO)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "Tax Engine · The 7-Step Calculation", size=28, bold=True, color=WHITE)

    steps = [
        ("Step 1", "Federal Taxable Income", "max(0, GrossIncome − Deductions)"),
        ("Step 2", "QBI Deduction", "Pass-throughs only: 20% × FedTaxable. C-Corps skip."),
        ("Step 3", "NOL Carryforward", "min(NOL, 80% × AdjustedTaxable)"),
        ("Step 4", "Federal Tax", "C-Corp flat 21% · Pass-throughs use 2024 brackets"),
        ("Step 5", "Apply Credits", "max(0, FedTax − Credits) — never negative"),
        ("Step 6", "State Tax", "Pre-QBI base × state rate · CA min $800 franchise"),
        ("Step 7", "Total + Effective Rate", "Total = Federal + State · Eff = Total/Gross"),
    ]
    cols = 4; rows = 2
    cw, rh = Inches(3.0), Inches(2.5)
    sx, sy = Inches(0.5), Inches(1.2)
    for i, (n, name, formula) in enumerate(steps):
        cx = sx + (i % cols) * (cw + Inches(0.15))
        cy = sy + (i // cols) * (rh + Inches(0.15))
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, rh)
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = INDIGO; card.line.width = Pt(1.25)
        add_text(s, cx + Inches(0.15), cy + Inches(0.15), cw, Inches(0.4),
                 n, size=12, bold=True, color=INDIGO)
        add_text(s, cx + Inches(0.15), cy + Inches(0.55), cw - Inches(0.3), Inches(0.5),
                 name, size=15, bold=True, color=SLATE)
        add_text(s, cx + Inches(0.15), cy + Inches(1.1), cw - Inches(0.3), Inches(1.3),
                 formula, size=11, color=GRAY)

    add_footer(s, 4, TOTAL)

# ==============================================================================
# SLIDE 5 — 13 AGENT TOOLS
# ==============================================================================
def slide_tools():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, WHITE)
    add_band(s, 0, Inches(0.9), INDIGO)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "13 Agent Tools · Function Calling", size=28, bold=True, color=WHITE)

    tools = [
        ("get_company_tax_breakdown", "Full 7-step breakdown"),
        ("what_if_calculator", "Override any field; see delta"),
        ("recalculate_with_overrides", "Multi-override scenarios"),
        ("compare_companies", "Side-by-side, 2+ companies"),
        ("portfolio_analysis", "Cross-portfolio insights"),
        ("filter_companies", "By entity / state / income / tax"),
        ("explain_tax_rule", "QBI · NOL · CA min · brackets"),
        ("nol_impact_analysis", "How much NOL is saving you"),
        ("credit_impact_analysis", "Tax saved by credits"),
        ("state_comparison", "Best state for this company"),
        ("entity_type_comparison", "Best entity type"),
        ("top_n_analysis", "Top N by any metric"),
        ("tax_savings_opportunities", "Restructuring picks"),
    ]
    cols = 3
    cw, rh = Inches(4.1), Inches(0.85)
    sx, sy = Inches(0.4), Inches(1.15)
    for i, (name, desc) in enumerate(tools):
        cx = sx + (i % cols) * (cw + Inches(0.15))
        cy = sy + (i // cols) * (rh + Inches(0.12))
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, rh)
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xC7, 0xD2, 0xFE); card.line.width = Pt(0.75)
        # left badge
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.1), cy + Inches(0.18), Inches(0.5), Inches(0.5))
        badge.fill.solid(); badge.fill.fore_color.rgb = INDIGO
        badge.line.fill.background()
        add_text(s, cx + Inches(0.1), cy + Inches(0.23), Inches(0.5), Inches(0.4),
                 str(i + 1), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, cx + Inches(0.7), cy + Inches(0.08), cw - Inches(0.8), Inches(0.4),
                 name, size=11, bold=True, color=INDIGO,
                 font="Consolas")
        add_text(s, cx + Inches(0.7), cy + Inches(0.42), cw - Inches(0.8), Inches(0.4),
                 desc, size=10, color=SLATE)

    add_footer(s, 5, TOTAL)

# ==============================================================================
# SLIDE 6 — DATA FLOW (CHAT TURN)
# ==============================================================================
def slide_flow():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, WHITE)
    add_band(s, 0, Inches(0.9), INDIGO)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "Data Flow · A Single Chat Turn", size=28, bold=True, color=WHITE)

    steps = [
        ("USER", "\"What if Pacific Manufacturing moved to Texas?\""),
        ("ANGULAR", "POST /api/chat { session_id, message }"),
        ("FASTAPI", "Load prior messages from MongoDB"),
        ("AI LAYER", "Send messages + 13 tool schemas to Azure OpenAI"),
        ("AZURE", "Returns tool_call: what_if_calculator(StateCode='TX')"),
        ("DISPATCHER", "Run tax_engine.what_if_calculator (deterministic)"),
        ("AZURE", "Receives JSON, composes natural-language reply"),
        ("MONGODB", "Persist user + assistant + tool messages"),
        ("ANGULAR", "Render reply + tool-call chips → user sees agent at work"),
    ]

    sy = Inches(1.2)
    rh = Inches(0.55)
    for i, (label, body) in enumerate(steps):
        y = sy + i * rh
        # Number circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y + Inches(0.05), Inches(0.45), Inches(0.45))
        circ.fill.solid(); circ.fill.fore_color.rgb = INDIGO
        circ.line.fill.background()
        add_text(s, Inches(0.5), y + Inches(0.1), Inches(0.45), Inches(0.4),
                 str(i + 1), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Label pill
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), y + Inches(0.05), Inches(1.6), Inches(0.45))
        pill.fill.solid(); pill.fill.fore_color.rgb = PURPLE
        pill.line.fill.background()
        add_text(s, Inches(1.1), y + Inches(0.13), Inches(1.6), Inches(0.4),
                 label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Body
        add_text(s, Inches(2.85), y + Inches(0.1), Inches(10), Inches(0.4),
                 body, size=14, color=SLATE)

    add_footer(s, 6, TOTAL)

# ==============================================================================
# SLIDE 7 — BACKEND FILES
# ==============================================================================
def slide_backend_files():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, LIGHT)
    add_band(s, 0, Inches(0.9), GREEN)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "Backend Files · backend/", size=28, bold=True, color=WHITE)

    files = [
        ("main.py", "FastAPI app + lifespan + all REST endpoints (companies, tax, chat, sessions)."),
        ("tax_engine.py", "Pure deterministic 7-step calculator + 13 tool functions + Excel loader."),
        ("ai_layer.py", "Azure OpenAI client, system prompt, 13 tool schemas, agentic loop with tool dispatch."),
        ("db.py", "MongoDB session store; transparent in-memory fallback if Mongo unavailable."),
        ("models.py", "Pydantic request/response models (Company, TaxResult, ChatRequest, …)."),
        ("requirements.txt", "fastapi · uvicorn · pydantic · pandas · openpyxl · openai · pymongo · dotenv"),
        (".env.example", "Template for AZURE_OPENAI_*, MONGO_URI, EXCEL_PATH."),
    ]
    sy = Inches(1.2)
    for i, (fname, desc) in enumerate(files):
        y = sy + i * Inches(0.78)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(12.4), Inches(0.65))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xD1, 0xFA, 0xE5); card.line.width = Pt(0.75)
        add_text(s, Inches(0.7), y + Inches(0.13), Inches(2.6), Inches(0.4),
                 fname, size=13, bold=True, color=GREEN, font="Consolas")
        add_text(s, Inches(3.4), y + Inches(0.13), Inches(9.5), Inches(0.4),
                 desc, size=12, color=SLATE)

    add_footer(s, 7, TOTAL)

# ==============================================================================
# SLIDE 8 — FRONTEND FILES
# ==============================================================================
def slide_frontend_files():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, LIGHT)
    add_band(s, 0, Inches(0.9), INDIGO)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "Frontend Files · frontend/src/", size=28, bold=True, color=WHITE)

    files = [
        ("app/tax-widget.component.ts", "Main UI: PrimeNG p-table grid + tax breakdown + chat panel + tool-chips."),
        ("app/api.service.ts", "Single HTTP service for all backend calls (companies, tax, chat, sessions)."),
        ("app/session.service.ts", "Reactive chat state via Angular signals; stable session_id per tab."),
        ("app/tax.types.ts", "TypeScript interfaces matching backend Pydantic models."),
        ("app/app.component.ts", "Shell — renders <tax-widget />."),
        ("app/app.config.ts", "Bootstrap providers: HttpClient + Animations."),
        ("main.ts · index.html · styles.css", "Standard Angular bootstrap + PrimeNG theme imports."),
        ("proxy.conf.json", "Dev proxy: /api → http://localhost:8000."),
    ]
    sy = Inches(1.2)
    for i, (fname, desc) in enumerate(files):
        y = sy + i * Inches(0.7)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(12.4), Inches(0.6))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xC7, 0xD2, 0xFE); card.line.width = Pt(0.75)
        add_text(s, Inches(0.7), y + Inches(0.1), Inches(3.7), Inches(0.4),
                 fname, size=12, bold=True, color=INDIGO, font="Consolas")
        add_text(s, Inches(4.5), y + Inches(0.1), Inches(8.4), Inches(0.4),
                 desc, size=12, color=SLATE)

    add_footer(s, 8, TOTAL)

# ==============================================================================
# SLIDE 9 — RUN STEPS
# ==============================================================================
def slide_run():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, WHITE)
    add_band(s, 0, Inches(0.9), INDIGO)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "How to Run · Three Terminals", size=28, bold=True, color=WHITE)

    # Box 1 — Mongo
    add_card(s, Inches(0.4), Inches(1.2), Inches(4.0), Inches(5.5),
             "1. MONGODB (optional)",
             [
                 "Local: mongodb://localhost:27017",
                 "Or Atlas SRV connection string",
                 "Skipping it? Backend uses in-memory fallback.",
                 "Set MONGO_URI in backend/.env",
             ], accent=GREEN)

    # Box 2 — Backend
    add_card(s, Inches(4.6), Inches(1.2), Inches(4.0), Inches(5.5),
             "2. BACKEND",
             [
                 "cd backend",
                 "python -m venv .venv",
                 ".\\.venv\\Scripts\\Activate.ps1",
                 "pip install -r requirements.txt",
                 "copy .env.example .env",
                 "→ fill in Azure + Mongo creds",
                 "uvicorn main:app --reload --port 8000",
                 "Verify: http://localhost:8000/companies",
             ], accent=ORANGE)

    # Box 3 — Frontend
    add_card(s, Inches(8.8), Inches(1.2), Inches(4.0), Inches(5.5),
             "3. FRONTEND",
             [
                 "cd frontend",
                 "npm install",
                 "npm start",
                 "Open http://localhost:4200",
                 "Click any company → see breakdown",
                 "Type in chat → AI agent responds",
                 "Tool-call chips show which tool fired",
             ], accent=INDIGO)

    add_footer(s, 9, TOTAL)

# ==============================================================================
# SLIDE 10 — DEMO SCRIPT
# ==============================================================================
def slide_demo():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, LIGHT)
    add_band(s, 0, Inches(0.9), INDIGO)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "Demo Script · 6 Power Prompts", size=28, bold=True, color=WHITE)

    demos = [
        ("Single calc",        "What is the total tax for Pacific Manufacturing Co?",
                                "→ get_company_tax_breakdown"),
        ("Step-by-step",       "Walk me through how Atlas Construction's federal tax was calculated.",
                                "→ get_company_tax_breakdown + narration"),
        ("What-if (state)",    "What if Pine Grove moved to Texas?",
                                "→ what_if_calculator"),
        ("Multi-override",     "What if Pine Grove was an LLC in CA with 200k more deductions?",
                                "→ recalculate_with_overrides"),
        ("Comparison",         "Compare Pacific Manufacturing and Atlas Construction.",
                                "→ compare_companies"),
        ("Portfolio insight",  "Which companies could save the most by restructuring?",
                                "→ tax_savings_opportunities"),
    ]
    sy = Inches(1.2)
    for i, (label, prompt, tool) in enumerate(demos):
        y = sy + i * Inches(0.95)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(12.4), Inches(0.85))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xC7, 0xD2, 0xFE)
        # Pill label
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y + Inches(0.18), Inches(2.0), Inches(0.5))
        pill.fill.solid(); pill.fill.fore_color.rgb = INDIGO
        pill.line.fill.background()
        add_text(s, Inches(0.7), y + Inches(0.27), Inches(2.0), Inches(0.4),
                 label, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.85), y + Inches(0.1), Inches(10), Inches(0.4),
                 prompt, size=12, color=SLATE)
        add_text(s, Inches(2.85), y + Inches(0.45), Inches(10), Inches(0.4),
                 tool, size=11, color=GREEN, font="Consolas")

    add_footer(s, 10, TOTAL)

# ==============================================================================
# SLIDE 11 — WHY WE WIN
# ==============================================================================
def slide_why():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, WHITE)
    add_band(s, 0, Inches(0.9), INDIGO)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "Why Tax Lens Wins", size=28, bold=True, color=WHITE)

    points = [
        ("Deterministic", "Tax math is pure Python. Auditable. Reproducible. No hallucinations.", INDIGO),
        ("Agentic", "13 callable tools — true function calling, not prompt-engineering tricks.", PURPLE),
        ("Conversational", "Real back-and-forth. Follow-ups. Challenges. Multi-part questions.", GREEN),
        ("Persistent", "MongoDB session memory — chat state survives restarts.", ORANGE),
        ("Transparent", "Every reply shows which tools fired — no black box.", INDIGO),
        ("Strategic", "Identifies restructuring opportunities across the portfolio.", PURPLE),
    ]
    cw = Inches(4.1); rh = Inches(2.6)
    sx = Inches(0.4); sy = Inches(1.2)
    for i, (title, desc, color) in enumerate(points):
        cx = sx + (i % 3) * (cw + Inches(0.15))
        cy = sy + (i // 3) * (rh + Inches(0.2))
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, rh)
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = color; card.line.width = Pt(2)
        add_text(s, cx + Inches(0.25), cy + Inches(0.25), cw - Inches(0.5), Inches(0.6),
                 title, size=22, bold=True, color=color)
        add_text(s, cx + Inches(0.25), cy + Inches(1.0), cw - Inches(0.5), Inches(1.5),
                 desc, size=13, color=SLATE)

    add_footer(s, 11, TOTAL)

# ==============================================================================
# SLIDE 12 — THANK YOU
# ==============================================================================
def slide_thanks():
    s = prs.slides.add_slide(BLANK)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    band.fill.solid(); band.fill.fore_color.rgb = INDIGO
    band.line.fill.background()

    add_text(s, Inches(0.8), Inches(2.8), Inches(11), Inches(1),
             "Thank You", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(4.2), Inches(11), Inches(0.6),
             "Tax Lens · Built in 2 hours · Powered by deterministic math + agentic AI",
             size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5),
             "Questions?", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


SLIDES = [
    slide_title,
    slide_problem,
    slide_architecture,
    slide_tax_steps,
    slide_tools,
    slide_flow,
    slide_backend_files,
    slide_frontend_files,
    slide_run,
    slide_demo,
    slide_why,
    slide_thanks,
]
TOTAL = len(SLIDES)

for fn in SLIDES:
    fn()

prs.save("presentation.pptx")
print(f"Generated presentation.pptx with {TOTAL} slides")
